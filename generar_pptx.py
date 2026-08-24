# -*- coding: utf-8 -*-
"""Genera una presentación de PowerPoint didáctica sobre el código ArraysUtils.java."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Paleta de colores (tema oscuro de programación) ---
BG = RGBColor(0x1E, 0x1E, 0x2E)          # fondo general
CARD = RGBColor(0x28, 0x28, 0x3A)        # tarjetas / cajas
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)      # azul claro
GREEN = RGBColor(0x3F, 0xD9, 0x8B)       # verde
YELLOW = RGBColor(0xFF, 0xE0, 0x6B)      # amarillo
RED = RGBColor(0xFF, 0x6B, 0x6B)         # rojo
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0xB0, 0xB8, 0xC8)
CODE_BG = RGBColor(0x11, 0x11, 0x1C)     # fondo del código

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FONT = "Consolas"
BODY_FONT = "Segoe UI"


def slide():
    s = prs.slides.add_slide(BLANK)
    # fondo
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, color=CARD, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    if line:
        r.line.color.rgb = line
        r.line.width = Pt(1.2)
    r.shadow.inherit = False
    r.adjustments[0] = 0.06
    return r


def txt(s, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        font=BODY_FONT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return tb


def title(s, text, subtitle=None):
    box(s, 0, 0, SW, Inches(1.15), ACCENT)
    txt(s, Inches(0.6), Inches(0.18), SW - Inches(1.2), Inches(0.8), text,
        size=32, color=RGBColor(0x0E, 0x1B, 0x2A), bold=True)
    if subtitle:
        txt(s, Inches(0.6), Inches(1.25), SW - Inches(1.2), Inches(0.5), subtitle,
            size=16, color=GRAY)


def footer(s, n):
    txt(s, SW - Inches(1.2), SH - Inches(0.5), Inches(1), Inches(0.4), str(n),
        size=14, color=GRAY, align=PP_ALIGN.RIGHT)


# ============================ DIAPOSITIVA 1: PORTADA ============================
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(1.1), SW - Inches(1.6), Inches(1.0),
    "Arreglos (Arrays) en Java", size=52, bold=True, color=RGBColor(0x0E, 0x1B, 0x2A))
txt(s, Inches(0.8), Inches(2.1), SW - Inches(1.6), Inches(0.8),
    "Análisis del programa ArraysUtils", size=28, color=RGBColor(0x0E, 0x1B, 0x2A))
box(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(2.6), CARD)
txt(s, Inches(1.1), Inches(3.4), Inches(11), Inches(2.2),
    "Lo que aprenderás:\n"
    "\u2022  Declarar y recorrer arreglos de enteros\n"
    "\u2022  Leer datos por teclado con Scanner\n"
    "\u2022  Recorrer con el bucle for-each\n"
    "\u2022  Contar pares, hallar máximo y mínimo\n"
    "\u2022  Invertir el orden de un arreglo",
    size=20, color=WHITE)
footer(s, 1)

# ============================ DIAPOSITIVA 2: VISIÓN GENERAL ============================
s = slide()
title(s, "¿Qué hace este programa?", "Estructura general")
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(1.0),
    "El programa lee 8 números por teclado, los guarda en un arreglo y calcula "
    "distintos datos sobre ellos.", size=19, color=WHITE)
pasos = [
    ("1. Leer 8 números", "Scanner + un bucle for los almacena en nums[]"),
    ("2. Contar pares", "¿Cuántos números son divisibles entre 2?"),
    ("3. Mayor y menor", "Busca el valor máximo y el mínimo del arreglo"),
    ("4. Invertir", "Crea un nuevo arreglo con el orden contrario"),
]
y = Inches(3.0)
for t, d in pasos:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.85), CARD)
    txt(s, Inches(1.0), y + Inches(0.12), Inches(4.4), Inches(0.6), t,
        size=20, color=ACCENT, bold=True)
    txt(s, Inches(5.6), y + Inches(0.16), Inches(6.8), Inches(0.6), d, size=16, color=GRAY)
    y += Inches(1.0)
footer(s, 2)

# ============================ DIAPOSITIVA 3: EL CÓDIGO COMPLETO ============================
s = slide()
title(s, "El código completo", "ArraysUtils.java")
code = (
    "import java.util.Scanner;\n"
    "\n"
    "public class ArraysUtils {\n"
    "    public static void main(String[] args) {\n"
    "        Scanner scanner = new Scanner(System.in);\n"
    "        int[] nums = new int[8];\n"
    "\n"
    "        for (int indice = 0; indice < nums.length; indice++) {\n"
    "            nums[indice] = scanner.nextInt();\n"
    "        }\n"
    "\n"
    "        System.out.println(\"Pares: \" + contarPares(nums));\n"
    "        System.out.println(\"Mayor: \" + mayor(nums)\n"
    "                           + \" Menor: \" + menor(nums));\n"
    "\n"
    "        int[] invertidos = invertir(nums);\n"
    "        System.out.println(\"Invertido: \"\n"
    "                           + java.util.Arrays.toString(invertidos));\n"
    "    }\n"
    "    // ... métodos contarPares, mayor, menor, invertir\n"
    "}"
)
box(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.1), CODE_BG, line=ACCENT)
txt(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.7), code, size=15, color=WHITE,
    font=FONT)
footer(s, 3)

# ============================ DIAPOSITIVA 4: DECLARACIÓN Y LECTURA ============================
s = slide()
title(s, "Declaración y lectura de datos", "Parte 1")
txt(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.7),
    "Declarar el arreglo:", size=20, color=ACCENT, bold=True)
box(s, Inches(0.7), Inches(2.3), Inches(5.8), Inches(1.0), CODE_BG)
txt(s, Inches(0.9), Inches(2.5), Inches(5.4), Inches(0.7),
    "int[] nums = new int[8];", size=18, color=WHITE, font=FONT)
txt(s, Inches(0.7), Inches(3.5), Inches(5.8), Inches(1.6),
    "\u2022  Crea un arreglo de 8 enteros.\n"
    "\u2022  Los índices van de 0 a 7.\n"
    "\u2022  nums.length vale 8.", size=16, color=GRAY)

txt(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(0.7),
    "Leer con Scanner:", size=20, color=ACCENT, bold=True)
box(s, Inches(6.9), Inches(2.3), Inches(5.8), Inches(1.5), CODE_BG)
txt(s, Inches(7.1), Inches(2.5), Inches(5.4), Inches(1.2),
    "for (int i = 0; i < nums.length; i++) {\n"
    "    nums[i] = scanner.nextInt();\n"
    "}", size=16, color=WHITE, font=FONT)
txt(s, Inches(6.9), Inches(4.0), Inches(5.8), Inches(1.6),
    "\u2022  Repite 8 veces pidiendo un entero.\n"
    "\u2022  Cada valor se guarda en una posición.\n"
    "\u2022  Scanner lee lo que escribimos por teclado.",
    size=16, color=GRAY)
footer(s, 4)

# ============================ DIAPOSITIVA 5: RECORRIDO FOR-EACH ============================
s = slide()
title(s, "El bucle for-each", "Recorrer sin usar índices")
box(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.3), CODE_BG)
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.0),
    "for (int valor : numeros) {  ...  }", size=20, color=WHITE, font=FONT)
txt(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(2.2),
    "\u2022  Recorre cada elemento del arreglo de principio a fin.\n"
    "\u2022  En cada vuelta, la variable valor toma el elemento actual.\n"
    "\u2022  Es más simple que el bucle for con índice cuando no necesitamos "
    "la posición.\n"
    "\u2022  Se usa en los métodos contarPares, mayor y menor.",
    size=18, color=WHITE)
box(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.0), CARD)
txt(s, Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.7),
    "Contraste: el bucle for con índice sí es necesario en invertir, porque "
    "ahí sí importa la posición.", size=16, color=YELLOW)
footer(s, 5)

# ============================ DIAPOSITIVA 6: CONTAR PARES ============================
s = slide()
title(s, "contarPares", "¿Cuántos números son pares?")
box(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(2.6), CODE_BG)
txt(s, Inches(0.9), Inches(1.9), Inches(5.6), Inches(2.3),
    "static int contarPares(int[] numeros) {\n"
    "    int contadorPares = 0;\n"
    "    for (int valor : numeros) {\n"
    "        if (valor % 2 == 0) {\n"
    "            contadorPares++;\n"
    "        }\n"
    "    }\n"
    "    return contadorPares;\n"
    "}", size=14, color=WHITE, font=FONT)
txt(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.5),
    "Paso a paso:", size=18, color=ACCENT, bold=True)
explicacion = [
    "1. Inicia un contador en 0.",
    "2. Recorre cada valor del arreglo.",
    "3. Si valor % 2 == 0, el número es par → suma 1.",
    "4. Devuelve el total de pares.",
]
y = Inches(2.6)
for e in explicacion:
    txt(s, Inches(7.0), y, Inches(5.6), Inches(0.9), e, size=16, color=WHITE)
    y += Inches(0.95)
txt(s, Inches(0.7), Inches(4.6), Inches(6.0), Inches(0.8),
    "Ejemplo: [1, 2, 3, 4]  →  resultado 2", size=16, color=GREEN, bold=True)
footer(s, 6)

# ============================ DIAPOSITIVA 7: MAYOR Y MENOR ============================
s = slide()
title(s, "mayor y menor", "Hallar el máximo y el mínimo")
box(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(3.1), CODE_BG)
txt(s, Inches(0.9), Inches(1.9), Inches(5.6), Inches(2.8),
    "static int mayor(int[] numeros) {\n"
    "    int maximo = numeros[0];\n"
    "    for (int valor : numeros) {\n"
    "        if (valor > maximo) {\n"
    "            maximo = valor;\n"
    "        }\n"
    "    }\n"
    "    return maximo;\n"
    "}", size=14, color=WHITE, font=FONT)
txt(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.3),
    "Idea clave (mismo patrón en ambos):", size=18, color=ACCENT, bold=True)
explicacion = [
    "1. Tomamos el primer valor como candidato.",
    "2. Comparamos con cada elemento.",
    "3. Si encontramos uno mayor, lo actualizamos.",
    "4. menor funciona igual, pero buscando valores menores.",
]
y = Inches(2.6)
for e in explicacion:
    txt(s, Inches(7.0), y, Inches(5.6), Inches(0.95), e, size=16, color=WHITE)
    y += Inches(0.95)
txt(s, Inches(0.7), Inches(5.1), Inches(6.0), Inches(1.0),
    "Ejemplo: [5, 9, 2, 7]  →  mayor 9, menor 2", size=16, color=GREEN, bold=True)
footer(s, 7)

# ============================ DIAPOSITIVA 8: INVERTIR ============================
s = slide()
title(s, "invertir", "Dar la vuelta al arreglo")
box(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(2.9), CODE_BG)
txt(s, Inches(0.9), Inches(1.9), Inches(5.6), Inches(2.6),
    "static int[] invertir(int[] numeros) {\n"
    "    int[] invertido =\n"
    "            new int[numeros.length];\n"
    "    for (int i = 0; i < numeros.length; i++) {\n"
    "        invertido[i] =\n"
    "           numeros[numeros.length - 1 - i];\n"
    "    }\n"
    "    return invertido;\n"
    "}", size=13, color=WHITE, font=FONT)
txt(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.3),
    "¿Cómo funciona?", size=18, color=ACCENT, bold=True)
explicacion = [
    "1. Se crea un arreglo nuevo del mismo tamaño.",
    "2. Recorremos con índice (aquí importa la posición).",
    "3. invertido[0] toma numeros[7],\n   invertido[1] → numeros[6], etc.",
    "4. La expresión length-1-i invierte el índice.",
]
y = Inches(2.6)
for e in explicacion:
    txt(s, Inches(7.0), y, Inches(5.6), Inches(1.1), e, size=16, color=WHITE)
    y += Inches(1.05)
txt(s, Inches(0.7), Inches(4.9), Inches(6.0), Inches(1.0),
    "Ejemplo: [1,2,3,4]  →  [4,3,2,1]", size=16, color=GREEN, bold=True)
footer(s, 8)

# ============================ DIAPOSITIVA 9: EJEMPLO COMPLETO ============================
s = slide()
title(s, "Ejemplo de ejecución", "Entrada y salida")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.6),
    "Si escribimos estos 8 números:", size=18, color=WHITE)
box(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(0.9), CODE_BG)
txt(s, Inches(0.9), Inches(2.5), Inches(5.2), Inches(0.6),
    "3  8  5  12  4  9  2  7", size=22, color=GREEN, font=FONT)
txt(s, Inches(6.7), Inches(1.7), Inches(6.0), Inches(0.6),
    "El programa mostrará:", size=18, color=WHITE)
box(s, Inches(6.7), Inches(2.3), Inches(6.0), Inches(2.0), CODE_BG)
txt(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(1.6),
    "Pares: 4\n"
    "Mayor: 12 Menor: 2\n"
    "Invertido: [7, 2, 9, 4, 12, 5, 8, 3]",
    size=17, color=WHITE, font=FONT)
box(s, Inches(0.7), Inches(3.6), Inches(5.6), Inches(2.6), CARD)
txt(s, Inches(1.0), Inches(3.8), Inches(5.0), Inches(2.2),
    "Comprobación:\n"
    "\u2022  Pares: 8, 12, 4, 2 → 4 ✓\n"
    "\u2022  Mayor: 12  |  Menor: 2 ✓\n"
    "\u2022  Invertido da la vuelta a la lista ✓",
    size=16, color=WHITE)
footer(s, 9)

# ============================ DIAPOSITIVA 10: RESUMEN ============================
s = slide()
title(s, "Resumen", "Ideas para recordar")
ideas = [
    ("Declaración", "int[] nums = new int[8]; crea un arreglo de 8 enteros."),
    ("Leer", "Un bucle for con scanner.nextInt() llena el arreglo."),
    ("for-each", "Recorre elementos sin necesitar el índice."),
    ("Contador", "contarPares usa una variable que suma cada vez que hay par."),
    ("Máx / mín", "Se parte del primer valor y se va actualizando al comparar."),
    ("Invertir", "Numeros.length-1-i devuelve la posición contraria."),
]
y = Inches(1.8)
for t, d in ideas:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.72), CARD)
    txt(s, Inches(1.0), y + Inches(0.1), Inches(3.4), Inches(0.5), t,
        size=18, color=ACCENT, bold=True)
    txt(s, Inches(4.6), y + Inches(0.12), Inches(7.8), Inches(0.5), d, size=15, color=WHITE)
    y += Inches(0.82)
footer(s, 10)

# ============================ DIAPOSITIVA 11: GRACIAS ============================
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(1.2),
    "¡Gracias!", size=60, bold=True, color=RGBColor(0x0E, 0x1B, 0x2A),
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.8),
    "¿Alguna pregunta?", size=28, color=RGBColor(0x0E, 0x1B, 0x2A),
    align=PP_ALIGN.CENTER)
footer(s, 11)

out = r"C:\Users\Pablo\Desktop\Pablo\anki-cards-completo\ArraysUtils_Java.pptx"
prs.save(out)
print("Guardado en:", out)
