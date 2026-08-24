# -*- coding: utf-8 -*-
"""Genera presentación de repaso para EXAMEN ESCRITO en papel (Arrays en Java)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Tema oscuro estilo "exámenes" (morado #7C6CF0) + cajas y emojis (gusto de Pablo)
BG = RGBColor(0x1A, 0x1B, 0x2E)
CARD = RGBColor(0x25, 0x26, 0x3D)
CARD2 = RGBColor(0x2E, 0x2F, 0x4A)
ACCENT = RGBColor(0x7C, 0x6C, 0xF0)   # morado exámenes
PURPLE_L = RGBColor(0xA9, 0x9C, 0xFF)
GREEN = RGBColor(0x34, 0xA8, 0x53)
YELLOW = RGBColor(0xFF, 0xD7, 0x66)
RED = RGBColor(0xEA, 0x43, 0x35)
BLUE = RGBColor(0x42, 0x85, 0xF4)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0xB0, 0xB8, 0xC8)
CODE_BG = RGBColor(0x0E, 0x0F, 0x1A)

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FONT = "Consolas"
BF = "Segoe UI"


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, color=CARD, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    if line:
        r.line.color.rgb = line; r.line.width = Pt(1.5)
    r.shadow.inherit = False; r.adjustments[0] = 0.08
    return r


def txt(s, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        font=BF, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font
    return tb


def title(s, text, subtitle=None):
    box(s, 0, 0, SW, Inches(1.15), ACCENT)
    txt(s, Inches(0.6), Inches(0.16), SW - Inches(1.2), Inches(0.85), text,
        size=30, color=RGBColor(0x12, 0x0F, 0x24), bold=True)
    if subtitle:
        txt(s, Inches(0.6), Inches(1.25), SW - Inches(1.2), Inches(0.5), subtitle,
            size=16, color=GRAY)


def footer(s, n):
    txt(s, SW - Inches(1.1), SH - Inches(0.5), Inches(1), Inches(0.4), str(n),
        size=14, color=GRAY, align=PP_ALIGN.RIGHT)


def code(s, x, y, w, h, text, size=15, line=ACCENT):
    box(s, x, y, w, h, CODE_BG, line=line)
    txt(s, x + Inches(0.25), y + Inches(0.15), w - Inches(0.5), h - Inches(0.3),
        text, size=size, color=WHITE, font=FONT)


# ============ 1. PORTADA ============
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(1.0),
    "📝 REPASO PARA EXAMEN EN PAPEL", size=44, bold=True,
    color=RGBColor(0x12, 0x0F, 0x24))
txt(s, Inches(0.8), Inches(2.0), SW - Inches(1.6), Inches(0.8),
    "Arreglos (Arrays) en Java — ArraysUtils", size=30,
    color=RGBColor(0x12, 0x0F, 0x24))
box(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.9), RGBColor(0x33, 0x2A, 0x59))
txt(s, Inches(1.1), Inches(3.3), Inches(11.1), Inches(2.5),
    "🎯 Objetivo: aprobar escribiendo código a mano\n\n"
    "✅ Qué tienes que saber hacer SIN computadora:\n"
    "   • Declarar y llenar un arreglo\n"
    "   • Recorrerlo con for y con for-each\n"
    "   • Contar pares · hallar máximo y mínimo\n"
    "   • Invertir el arreglo\n"
    "   • Predecir la salida (trazar) y detectar errores",
    size=18, color=WHITE)
footer(s, 1)

# ============ 2. QUÉ TE PUEDEN PREGUNTAR ============
s = slide()
title(s, "🤔 Qué te pueden preguntar", "Tipos típicos de pregunta en examen escrito")
items = [
    ("A", "Completar el código", "Rellenar huecos en un método (líneas con ___).", BLUE),
    ("B", "Predecir la salida", "Dado un arreglo, escribir qué imprime el programa.", GREEN),
    ("C", "Escribir un método", "Implementar contarPares, mayor, menor o invertir desde cero.", YELLOW),
    ("D", "Trazar (trace)", "Hacer una tabla con los valores de las variables en cada vuelta.", PURPLE_L),
    ("E", "Detectar errores", "Señalar bugs (fuera de rango, índices, paridad).", RED),
]
y = Inches(1.75)
for letra, t, d, c in items:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.92), CARD)
    bx = box(s, Inches(0.95), y + Inches(0.17), Inches(0.6), Inches(0.58), c)
    txt(s, Inches(1.15), y + Inches(0.27), Inches(0.3), Inches(0.4), letra,
        size=20, bold=True, color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
    txt(s, Inches(1.85), y + Inches(0.13), Inches(4.0), Inches(0.7), t,
        size=20, bold=True, color=WHITE)
    txt(s, Inches(6.1), y + Inches(0.17), Inches(6.3), Inches(0.65), d,
        size=15, color=GRAY)
    y += Inches(1.03)
footer(s, 2)

# ============ 3. SINTAXIS A MEMORIZAR ============
s = slide()
title(s, "🧠 Sintaxis imprescindible", "Para escribir de memoria")
txt(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(0.5),
    "🟦 Declarar un arreglo", size=19, bold=True, color=BLUE)
code(s, Inches(0.7), Inches(2.2), Inches(5.9), Inches(1.15),
    "int[] nums = new int[8];", 17)
txt(s, Inches(0.7), Inches(3.45), Inches(5.9), Inches(0.9),
    "8 = cantidad de elementos. Índices 0..7.\nTamaño fijo, no cambia.",
    size=14, color=GRAY)

txt(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(0.5),
    "🟩 Llenar con Scanner", size=19, bold=True, color=GREEN)
code(s, Inches(6.9), Inches(2.2), Inches(5.9), Inches(1.5),
    "for (int i = 0; i < nums.length; i++) {\n"
    "    nums[i] = scanner.nextInt();\n"
    "}", 15)
txt(s, Inches(6.9), Inches(3.85), Inches(5.9), Inches(0.9),
    "Pide 8 enteros por teclado y los guarda.",
    size=14, color=GRAY)

txt(s, Inches(0.7), Inches(4.5), Inches(5.9), Inches(0.5),
    "🟨 Recorrer sin índice (for-each)", size=19, bold=True, color=YELLOW)
code(s, Inches(0.7), Inches(5.0), Inches(5.9), Inches(1.3),
    "for (int valor : numeros) {\n"
    "    // usa valor\n"
    "}", 15)

txt(s, Inches(6.9), Inches(4.5), Inches(5.9), Inches(0.5),
    "🟥 Dato clave: longitud", size=19, bold=True, color=RED)
code(s, Inches(6.9), Inches(5.0), Inches(5.9), Inches(1.3),
    "nums.length  // → 8\n"
    "último = nums[nums.length - 1];", 15)
footer(s, 3)

# ============ 4. TABLA DE TRAZADO ============
s = slide()
title(s, "📋 Cómo trazar a mano (trace)", "Haz una tabla en tu hoja de examen")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
    "Para el arreglo nums = [3, 8, 5, 12, 4, 9, 2, 7], traza contarPares:",
    size=16, color=WHITE)
# cabecera de tabla
cols = ["valor", "¿valor % 2 == 0?", "contadorPares"]
x = Inches(0.7); w = [Inches(2.2), Inches(5.4), Inches(4.3)]
for i, c in enumerate(cols):
    xpos = x + sum(w[j] for j in range(i))
    box(s, xpos, Inches(2.3), w[i], Inches(0.55), ACCENT)
    txt(s, xpos + Inches(0.15), Inches(2.35), w[i] - Inches(0.3), Inches(0.45),
        c, size=16, bold=True, color=RGBColor(0x12, 0x0F, 0x24))
rows = [
    ("3", "no", "0"), ("8", "sí ✔", "1"), ("5", "no", "1"), ("12", "sí ✔", "2"),
    ("4", "sí ✔", "3"), ("9", "no", "3"), ("2", "sí ✔", "4"), ("7", "no", "4"),
]
y = Inches(2.85)
for r in rows:
    for i, val in enumerate(r):
        xpos = x + sum(w[j] for j in range(i))
        box(s, xpos, y, w[i], Inches(0.5), CARD, line=CARD2)
        txt(s, xpos + Inches(0.15), y + Inches(0.06), w[i] - Inches(0.3),
            Inches(0.4), val, size=14, color=WHITE)
    y += Inches(0.5)
txt(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
    "🟩 Resultado: contarPares(nums) → 4   |  Poner así la tabla te quita errores.",
    size=17, bold=True, color=GREEN)
footer(s, 4)

# ============ 5. MAYOR Y MENOR (traza) ============
s = slide()
title(s, "🟣 Traza de mayor y menor", "Técnica del 'candidato'")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
    "Truco para el examen: el primero es el candidato inicial, luego se compara. "
    "Solo se actualiza si aparece un valor mejor.", size=16, color=WHITE)
code(s, Inches(0.7), Inches(2.35), Inches(6.3), Inches(2.2),
    "int maximo = numeros[0];   // 3\n"
    "for (int valor : numeros) {\n"
    "    if (valor > maximo) {\n"
    "        maximo = valor;\n"
    "    }\n"
    "}\n"
    "// maximo → 12", 14)
code(s, Inches(7.2), Inches(2.35), Inches(5.4), Inches(2.2),
    "int minimo = numeros[0];   // 3\n"
    "for (int valor : numeros) {\n"
    "    if (valor < minimo) {\n"
    "        minimo = valor;\n"
    "    }\n"
    "}\n"
    "// minimo → 2", 14)
box(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.3), CARD)
txt(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.0),
    "⚠️ Puntos que suelen penalizar en el examen:\n"
    "   • Inicializar en 0 en vez de numeros[0]  →  falla con números negativos\n"
    "   • Confundir > con <  →  mayor/menor al revés\n"
    "   • Olvidar el return al final del método",
    size=15, color=YELLOW)
footer(s, 5)

# ============ 6. INVERTIR (traza visual) ============
s = slide()
title(s, "🔁 invertir paso a paso", "length - 1 - i")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
    "Ejemplo nums = [3, 8, 5, 12, 4, 9, 2, 7] → invertido = [7, 2, 9, 4, 12, 5, 8, 3]",
    size=16, color=WHITE)
# mostrar fórmula
code(s, Inches(0.7), Inches(2.35), Inches(11.9), Inches(1.5),
    "for (int i = 0; i < numeros.length; i++) {\n"
    "    invertido[i] = numeros[numeros.length - 1 - i];\n"
    "}", 16)
box(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(2.0), CARD)
txt(s, Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.7),
    "🧮 Con i = 0 → numeros[7] (último) va a invertido[0]\n"
    "   Con i = 1 → numeros[6] → invertido[1]\n"
    "   ... y así hasta i = 7 → numeros[0] → invertido[7]\n\n"
    "✅ Repasa: length - 1 = 7, menos i para ir 'desde atrás hacia adelante'.",
    size=16, color=WHITE)
footer(s, 6)

# ============ 7. PLANTILLAS PARA ESCRIBIR A MANO ============
s = slide()
title(s, "✍️ Plantillas listas para copiar", "Estructura de cada método (memoriza)")
code(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(2.6),
    "// CONTAR PARES\n"
    "static int contarPares(int[] nums) {\n"
    "    int c = 0;\n"
    "    for (int v : nums)\n"
    "        if (v % 2 == 0) c++;\n"
    "    return c;\n"
    "}", 14)
code(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(2.6),
    "// INVERTIR (devuelve arreglo)\n"
    "static int[] invertir(int[] nums) {\n"
    "    int[] inv = new int[nums.length];\n"
    "    for (int i = 0; i < nums.length; i++)\n"
    "        inv[i] = nums[nums.length-1-i];\n"
    "    return inv;\n"
    "}", 14)
code(s, Inches(0.7), Inches(4.55), Inches(5.9), Inches(1.7),
    "// MAYOR\n"
    "static int mayor(int[] nums) {\n"
    "    int m = nums[0];\n"
    "    for (int v : nums)\n"
    "        if (v > m) m = v;\n"
    "    return m;\n"
    "}", 14)
code(s, Inches(6.9), Inches(4.55), Inches(5.9), Inches(1.7),
    "// MENOR (igual pero con <)\n"
    "static int menor(int[] nums) {\n"
    "    int m = nums[0];\n"
    "    for (int v : nums)\n"
    "        if (v < m) m = v;\n"
    "    return m;\n"
    "}", 14)
footer(s, 7)

# ============ 8. EJERCICIOS DE PRÁCTICA ============
s = slide()
title(s, "🏋️ Ejercicios para practicar a mano", "Hazlos en papel antes del examen")
ex = [
    ("1", "nums = [10, 15, 7, 22, 9]", "¿Qué devuelve contarPares? (respuesta: 3)"),
    ("2", "nums = [5, 2, 9, 1, 7]", "Escribe mayor y menor sin mirar. (9 y 1)"),
    ("3", "nums = [1, 2, 3, 4]", "Escribe la salida de invertir. ([4,3,2,1])"),
    ("4", "nums = [0, -5, 3, -2]", "¿Fallaría mayor si inicializas en 0? (sí → usaría nums[0])"),
    ("5", "Traza: nums = [4, 6, 3]", "Tabla completa de contarPares. (2)"),
]
y = Inches(1.8)
for num, e, d in ex:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.9), CARD)
    bx = box(s, Inches(0.95), y + Inches(0.18), Inches(0.55), Inches(0.55), ACCENT)
    txt(s, Inches(1.12), y + Inches(0.28), Inches(0.3), Inches(0.4), num,
        size=18, bold=True, color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
    txt(s, Inches(1.85), y + Inches(0.15), Inches(5.4), Inches(0.6), e,
        size=16, color=BLUE, font=FONT)
    txt(s, Inches(7.4), y + Inches(0.17), Inches(5.0), Inches(0.6), d,
        size=14, color=GRAY)
    y += Inches(1.0)
footer(s, 8)

# ============ 9. ERRORES COMUNES EN EXAMEN ============
s = slide()
title(s, "🚨 Errores típicos (no los cometas)", "Lo que más suele restar puntos")
errs = [
    ("Fuera de rango", "nums[nums.length] NO existe. El último es nums.length-1.", RED),
    ("Contar pares mal", "Usar valor % 2 == 1 cuenta IMPARES, no pares.", YELLOW),
    ("Inicializar máximo en 0", "Con negativos falla. Inicializa siempre en nums[0].", YELLOW),
    ("Confundir índices", "i++ vs i-- : en invertir, cuidado con length-1-i.", RED),
    ("Olvidar return", "Todo método int[] / int debe terminar con return.", RED),
    ("Tamaño fijo", "Los arreglos NO crecen solos: new int[8] es fijo.", YELLOW),
]
y = Inches(1.8)
for t, d, c in errs:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.75), CARD)
    box(s, Inches(0.95), y + Inches(0.15), Inches(0.12), Inches(0.45), c)
    txt(s, Inches(1.4), y + Inches(0.12), Inches(4.2), Inches(0.5), t,
        size=17, bold=True, color=WHITE)
    txt(s, Inches(5.8), y + Inches(0.15), Inches(6.6), Inches(0.5), d,
        size=14, color=GRAY)
    y += Inches(0.83)
footer(s, 9)

# ============ 10. CHECKLIST PRE-EXAMEN ============
s = slide()
title(s, "✅ Checklist antes de entrar", "¿Sabes hacer esto sin mirar?")
checks = [
    "Declarar int[] y asignar con new",
    "Llenar con Scanner en un for (índices 0..length-1)",
    "Recorrer con for-each",
    "Escribir contarPares (paridad con % 2)",
    "Escribir mayor y menor (candidato = nums[0])",
    "Escribir invertir (length-1-i)",
    "Trazar cualquier método en una tabla",
    "Leer la salida con println y Arrays.toString",
]
y = Inches(1.8)
for i, c in enumerate(checks):
    row = Inches(0.8)
    col = 0 if i < 4 else 1
    xx = Inches(0.7) + col * Inches(6.1)
    yy = y + (i % 4) * row
    box(s, xx, yy, Inches(5.9), Inches(0.7), CARD)
    box(s, xx + Inches(0.15), yy + Inches(0.16), Inches(0.4), Inches(0.38), GREEN)
    txt(s, xx + Inches(0.75), yy + Inches(0.14), Inches(5.0), Inches(0.5), c,
        size=15, color=WHITE)
txt(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.6),
    "🟣 Si marcas todos → estás listo para el examen. ¡Mucho ánimo! 💪",
    size=18, bold=True, color=PURPLE_L)
footer(s, 10)

# ============ 11. CIERRE ============
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(2.7), SW - Inches(1.6), Inches(1.1),
    "🍀 ¡Mucha suerte en el examen!", size=52, bold=True,
    color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.8),
    "Repasa las plantillas y practica los trazas a mano. Tú puedes. 💪",
    size=22, color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
footer(s, 11)

out = r"C:\Users\Pablo\Desktop\Pablo\anki-cards-completo\Arrays_Java_Repaso_Examen.pptx"
prs.save(out)
print("Guardado en:", out)
