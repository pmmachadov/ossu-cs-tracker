# -*- coding: utf-8 -*-
"""Presentación pedagógica con ilustraciones para examen escrito (Arrays en Java)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG = r"C:\Users\Pablo\Desktop\Pablo\anki-cards-completo\ilustraciones"

BG = RGBColor(0x14, 0x15, 0x27)
CARD = RGBColor(0x25, 0x26, 0x3D)
ACCENT = RGBColor(0x7C, 0x6C, 0xF0)
GREEN = RGBColor(0x34, 0xA8, 0x53)
YELLOW = RGBColor(0xFF, 0xD7, 0x66)
RED = RGBColor(0xEA, 0x43, 0x35)
BLUE = RGBColor(0x42, 0x85, 0xF4)
PURPLE_L = RGBColor(0xA9, 0x9C, 0xFF)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0xB0, 0xB8, 0xC8)
CODE_BG = RGBColor(0x0E, 0x0F, 0x1A)

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FONT = "Consolas"
BF = "Segoe UI"


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def box(s, x, y, w, h, color=CARD, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    if line:
        r.line.color.rgb = line; r.line.width = Pt(1.5)
    r.shadow.inherit = False; r.adjustments[0] = 0.08
    return r


def txt(s, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        font=BF, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font
    return tb


def title(s, text, subtitle=None):
    box(s, 0, 0, SW, Inches(1.15), ACCENT)
    txt(s, Inches(0.6), Inches(0.16), SW - Inches(1.2), Inches(0.85), text,
        size=30, color=RGBColor(0x12, 0x0F, 0x24), bold=True)
    if subtitle:
        txt(s, Inches(0.6), Inches(1.25), SW - Inches(1.2), Inches(0.5), subtitle,
            size=16, color=GRAY)


def footer(s, n):
    txt(s, SW - Inches(1.1), SH - Inches(0.5), Inches(1), Inches(0.4), str(n),
        size=14, color=GRAY, align=PP_ALIGN.RIGHT)


def code(s, x, y, w, h, text, size=15, line=ACCENT):
    box(s, x, y, w, h, CODE_BG, line=line)
    txt(s, x + Inches(0.25), y + Inches(0.15), w - Inches(0.5), h - Inches(0.3),
        text, size=size, color=WHITE, font=FONT)


def pic(s, name, x, y, w=None, h=None):
    p = os.path.join(IMG, name)
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    s.shapes.add_picture(p, x, y, **kw)


# ============ 1. PORTADA ============
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(0.9), SW - Inches(1.6), Inches(1.0),
    "📝 REPASO PARA EXAMEN EN PAPEL", size=44, bold=True,
    color=RGBColor(0x12, 0x0F, 0x24))
txt(s, Inches(0.8), Inches(1.9), SW - Inches(1.6), Inches(0.8),
    "Arreglos (Arrays) en Java — ArraysUtils", size=30,
    color=RGBColor(0x12, 0x0F, 0x24))
box(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(2.3), RGBColor(0x33, 0x2A, 0x59))
txt(s, Inches(1.1), Inches(3.1), Inches(11.1), Inches(1.9),
    "🎯 Escribe código a mano, SIN computadora\n\n"
    "✅ Declarar arreglos · llenar con Scanner · recorrer\n"
    "✅ Contar pares · máximo/mínimo · invertir\n"
    "✅ Predecir salidas (trazar) y detectar errores",
    size=19, color=WHITE)
footer(s, 1)

# ============ 2. QUÉ ES UN ARREGLO (con ilustración) ============
s = slide()
title(s, "🧱 ¿Qué es un arreglo?", "Una caja con varios valores en fila")
pic(s, "array_indices.png", Inches(0.7), Inches(1.9), w=Inches(11.9))
box(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.4), CARD)
txt(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.1),
    "🟪 Un arreglo guarda varios valores del mismo tipo bajo un solo nombre.\n"
    "🟩 Cada posición se llama con un índice que empieza en 0.\n"
    "🟦 El tamaño es fijo: new int[8] crea 8 huecos. El último hueco es [7].",
    size=16, color=WHITE)
footer(s, 2)

# ============ 3. DIAGRAMA DE FLUJO ============
s = slide()
title(s, "🗺️ Mapa del programa", "Secuencia de pasos (dibújalo en tu cabeza)")
pic(s, "flowchart.png", Inches(7.1), Inches(1.6), h=Inches(5.5))
box(s, Inches(0.7), Inches(1.9), Inches(6.1), Inches(4.4), CARD)
txt(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4.0),
    "🟩 1. INICIO\n\n"
    "🟪 2. Crear el arreglo\n     int[] nums = new int[8]\n\n"
    "🟦 3. Llenar en un bucle\n     8 números por teclado\n\n"
    "🟨 4. Procesar\n     contarPares · mayor · menor\n\n"
    "🟩 5. Invertir y mostrar\n\n"
    "🟥 6. FIN",
    size=18, color=WHITE)
footer(s, 3)

# ============ 4. DECLARAR Y LLENAR ============
s = slide()
title(s, "🛠️ Declarar y llenar", "Paso 1 y 2")
txt(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(0.5),
    "🟦 Declarar", size=19, bold=True, color=BLUE)
code(s, Inches(0.7), Inches(2.2), Inches(5.9), Inches(1.2),
    "int[] nums = new int[8];", 17)
txt(s, Inches(0.7), Inches(3.5), Inches(5.9), Inches(1.1),
    "Tipo int, nombre nums, 8 huecos.\nÍndices 0..7 · nums.length = 8",
    size=14, color=GRAY)

txt(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(0.5),
    "🟩 Llenar con Scanner", size=19, bold=True, color=GREEN)
code(s, Inches(6.9), Inches(2.2), Inches(5.9), Inches(1.6),
    "for (int i = 0; i < nums.length; i++) {\n"
    "    nums[i] = scanner.nextInt();\n"
    "}", 15)
txt(s, Inches(6.9), Inches(4.0), Inches(5.9), Inches(1.0),
    "Repite 8 veces pidiendo un entero\ny guardándolo en la posición i.",
    size=14, color=GRAY)

box(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.3), CARD)
txt(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.0),
    "🚨 En el examen: el bucle NUNCA debe llegar a nums[nums.length]. "
    "Con i < nums.length, i va de 0 a 7. La última posición válida es [7].",
    size=15, color=YELLOW)
footer(s, 4)

# ============ 5. for vs for-each (con ilustración) ============
s = slide()
title(s, "🔄 for vs for-each", "¿Cuándo usar cada uno?")
pic(s, "for_foreach.png", Inches(0.7), Inches(1.7), w=Inches(11.9))
footer(s, 5)

# ============ 6. contarPares (traza) ============
s = slide()
title(s, "🧮 contarPares", "Cuenta los pares paso a paso")
code(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(2.2),
    "static int contarPares(int[] n) {\n"
    "    int c = 0;\n"
    "    for (int v : n)\n"
    "        if (v % 2 == 0) c++;\n"
    "    return c;\n"
    "}", 14)
cols = ["valor", "¿%2==0?", "c"]
x = Inches(6.8); w = [Inches(2.0), Inches(2.6), Inches(1.4)]
for i, c in enumerate(cols):
    xpos = x + sum(w[j] for j in range(i))
    box(s, xpos, Inches(1.7), w[i], Inches(0.5), ACCENT)
    txt(s, xpos + Inches(0.1), Inches(1.75), w[i], Inches(0.4), c, size=14,
        bold=True, color=RGBColor(0x12, 0x0F, 0x24))
rows = [("3", "no", "0"), ("8", "sí", "1"), ("5", "no", "1"), ("12", "sí", "2"),
        ("4", "sí", "3"), ("9", "no", "3"), ("2", "sí", "4"), ("7", "no", "4")]
y = Inches(2.2)
for r in rows:
    for i, val in enumerate(r):
        xpos = x + sum(w[j] for j in range(i))
        box(s, xpos, y, w[i], Inches(0.48), CARD)
        txt(s, xpos + Inches(0.1), y + Inches(0.05), w[i], Inches(0.4), val,
            size=14, color=WHITE)
    y += Inches(0.48)
txt(s, Inches(0.7), Inches(4.2), Inches(5.8), Inches(0.9),
    "💡 v % 2 == 0 → par.\n💡 v % 2 == 1 → impar (¡trampa común!).",
    size=15, color=YELLOW)
box(s, Inches(6.8), Inches(6.2), Inches(5.8), Inches(0.8), GREEN)
txt(s, Inches(7.0), Inches(6.3), Inches(5.4), Inches(0.6),
    "Resultado: contarPares → 4", size=17, bold=True,
    color=RGBColor(0x12, 0x0F, 0x24))
footer(s, 6)

# ============ 7. mayor y menor (con ilustración) ============
s = slide()
title(s, "🟢 mayor y menor", "La técnica del 'candidato'")
pic(s, "mayor_menor.png", Inches(0.7), Inches(1.6), w=Inches(11.9))
box(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.4), CARD)
txt(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.1),
    "⚠️ Inicializar el máximo en nums[0] (no en 0): con números negativos "
    "inicializar en 0 fallaría. mayor usa > · menor usa <.",
    size=15, color=YELLOW)
footer(s, 7)

# ============ 8. invertir (con ilustración) ============
s = slide()
title(s, "🔁 invertir", "Dar la vuelta con length - 1 - i")
pic(s, "invert_trace.png", Inches(0.7), Inches(1.5), w=Inches(11.9))
footer(s, 8)

# ============ 9. PLANTILLAS ============
s = slide()
title(s, "✍️ Plantillas para escribir a mano", "Memoriza estas 4")
code(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(2.5),
    "// CONTAR PARES\n"
    "static int contarPares(int[] n) {\n"
    "    int c = 0;\n"
    "    for (int v : n)\n"
    "        if (v % 2 == 0) c++;\n"
    "    return c;\n"
    "}", 14)
code(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(2.5),
    "// INVERTIR\n"
    "static int[] invertir(int[] n) {\n"
    "    int[] inv = new int[n.length];\n"
    "    for (int i = 0; i < n.length; i++)\n"
    "        inv[i] = n[n.length-1-i];\n"
    "    return inv;\n"
    "}", 14)
code(s, Inches(0.7), Inches(4.45), Inches(5.9), Inches(1.6),
    "// MAYOR\n"
    "static int mayor(int[] n) {\n"
    "    int m = n[0];\n"
    "    for (int v : n)\n"
    "        if (v > m) m = v;\n"
    "    return m;\n"
    "}", 14)
code(s, Inches(6.9), Inches(4.45), Inches(5.9), Inches(1.6),
    "// MENOR (mismo con <)\n"
    "static int menor(int[] n) {\n"
    "    int m = n[0];\n"
    "    for (int v : n)\n"
    "        if (v < m) m = v;\n"
    "    return m;\n"
    "}", 14)
footer(s, 9)

# ============ 10. ERRORES COMUNES ============
s = slide()
title(s, "🚨 Errores que restan puntos", "No los cometas en el examen")
errs = [
    ("Fuera de rango", "nums[nums.length] NO existe. Último = length-1.", RED),
    ("Paridad al revés", "v % 2 == 1 cuenta IMPARES, no pares.", YELLOW),
    ("Máximo en 0", "Con negativos falla. Usa nums[0].", YELLOW),
    ("Índices en invertir", "length-1-i, con cuidado del -1.", RED),
    ("Olvidar return", "Todo método int/int[] termina con return.", RED),
    ("Tamaño fijo", "Un arreglo no crece: new int[8] es fijo.", YELLOW),
]
y = Inches(1.8)
for t, d, c in errs:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.72), CARD)
    box(s, Inches(0.95), y + Inches(0.14), Inches(0.12), Inches(0.44), c)
    txt(s, Inches(1.4), y + Inches(0.11), Inches(4.2), Inches(0.5), t,
        size=17, bold=True, color=WHITE)
    txt(s, Inches(5.8), y + Inches(0.14), Inches(6.6), Inches(0.5), d,
        size=14, color=GRAY)
    y += Inches(0.8)
footer(s, 10)

# ============ 11. EJERCICIOS ============
s = slide()
title(s, "🏋️ Practica en papel", "Hazlos antes del examen")
ex = [
    ("1", "nums = [10, 15, 7, 22, 9]", "contarPares → ? (3)"),
    ("2", "nums = [5, 2, 9, 1, 7]", "mayor y menor → ? (9 y 1)"),
    ("3", "nums = [1, 2, 3, 4]", "invertir → ? ([4,3,2,1])"),
    ("4", "nums = [0, -5, 3, -2]", "¿mayor falla si inicias en 0? (sí)"),
    ("5", "Traza nums = [4, 6, 3]", "tabla de contarPares → (2)"),
]
y = Inches(1.9)
for num, e, d in ex:
    box(s, Inches(0.7), y, Inches(11.9), Inches(0.85), CARD)
    bx = box(s, Inches(0.95), y + Inches(0.16), Inches(0.55), Inches(0.55), ACCENT)
    txt(s, Inches(1.12), y + Inches(0.26), Inches(0.3), Inches(0.4), num,
        size=18, bold=True, color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
    txt(s, Inches(1.85), y + Inches(0.15), Inches(5.3), Inches(0.6), e,
        size=16, color=BLUE, font=FONT)
    txt(s, Inches(7.3), y + Inches(0.17), Inches(5.0), Inches(0.6), d,
        size=14, color=GRAY)
    y += Inches(0.95)
box(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7), CARD)
txt(s, Inches(1.0), Inches(6.58), Inches(11.3), Inches(0.55),
    "🟣 Traza SIEMPRE en una tabla: te evita errores de cálculo en el examen.",
    size=16, bold=True, color=PURPLE_L)
footer(s, 11)

# ============ 12. CIERRE ============
s = slide()
box(s, 0, 0, SW, SH, ACCENT)
txt(s, Inches(0.8), Inches(2.7), SW - Inches(1.6), Inches(1.1),
    "🍀 ¡Mucha suerte en el examen!", size=52, bold=True,
    color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.8),
    "Repasa las plantillas y practica los trazas a mano. Tú puedes. 💪",
    size=22, color=RGBColor(0x12, 0x0F, 0x24), align=PP_ALIGN.CENTER)
footer(s, 12)

out = r"C:\Users\Pablo\Desktop\Pablo\anki-cards-completo\Arrays_Java_Pedagogico_Ilustrado.pptx"
prs.save(out)
print("Guardado en:", out)
